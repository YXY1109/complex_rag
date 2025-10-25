"""
Office Documents Handler

This module provides specialized processing for office documents including
PDF, DOCX, XLSX, PPTX and other office formats.
"""

import asyncio
import mimetypes
from typing import Dict, Any, Optional, List, Union
from pathlib import Path
from dataclasses import dataclass
from datetime import datetime
import io

try:
    import PyPDF2
    import fitz  # PyMuPDF
    PYPDF2_AVAILABLE = True
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from ...interfaces.source_interface import (
    SourceHandler,
    FileSource,
    ProcessingStrategy,
    ParseRequest,
    ParseResponse
)
from ...interfaces.parser_interface import (
    DocumentFormat,
    DocumentChunk,
    Metadata
)
from ...services.quality_monitor import QualityMonitor, QualityMetric


@dataclass
class OfficeDocumentFeatures:
    """Features extracted from office documents."""
    document_type: str
    page_count: int
    has_images: bool
    has_tables: bool
    has_formulas: bool
    slide_count: int
    worksheet_count: int
    section_count: int
    is_protected: bool
    has_comments: bool
    has_track_changes: bool


class OfficeDocumentsHandler(SourceHandler):
    """
    Handler for office documents including PDF, DOCX, XLSX, PPTX.

    Features:
    - Multi-format support (PDF, DOCX, XLSX, PPTX)
    - Text extraction with formatting preservation
    - Table extraction and processing
    - Image metadata extraction
    - Document structure analysis
    - Password protection handling
    """

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize office documents handler."""
        super().__init__(FileSource.OFFICE_DOCUMENTS, config)
        self.quality_monitor = QualityMonitor()
        self.supported_formats = self._get_supported_formats()

    def _get_supported_formats(self) -> Dict[str, str]:
        """Get supported office document formats."""
        formats = {
            '.pdf': 'pdf',
            '.doc': 'doc',
            '.docx': 'docx',
            '.xls': 'xls',
            '.xlsx': 'xlsx',
            '.ppt': 'ppt',
            '.pptx': 'pptx',
            '.odt': 'odt',
            '.ods': 'ods',
            '.odp': 'odp',
            '.rtf': 'rtf'
        }
        return formats

    async def connect(self) -> bool:
        """Initialize handler."""
        return True

    async def disconnect(self) -> None:
        """Cleanup handler."""
        pass

    async def can_handle(self, request: ParseRequest) -> bool:
        """Check if this handler can process the request."""
        # Check file extension
        if request.file_path:
            file_ext = Path(request.file_path).suffix.lower()
            if file_ext in self.supported_formats:
                return True

        # Check MIME type
        if request.mime_type:
            office_types = {
                'application/pdf',
                'application/msword',
                'application/vnd.ms-',
                'application/vnd.openxmlformats-',
                'application/vnd.oasis.opendocument.',
                'application/rtf'
            }
            if any(request.mime_type.startswith(prefix) for prefix in office_types):
                return True

        return False

    async def process(self, request: ParseRequest) -> ParseResponse:
        """Process office document."""
        session_id = f"office_{datetime.now().timestamp()}"
        processing_start = datetime.now()

        # Start quality monitoring
        quality_session = self.quality_monitor.start_session(
            session_id=session_id,
            file_source=FileSource.OFFICE_DOCUMENTS,
            strategy=request.strategy,
            file_size=len(request.content) if request.content else None
        )

        try:
            # Determine document format
            document_format = self._detect_office_format(request)
            document_type = self.supported_formats.get(document_format.value, 'unknown')

            # Extract content based on format
            if document_type == 'pdf':
                extracted_content = await self._process_pdf(request.content, request.custom_params)
            elif document_type == 'docx':
                extracted_content = await self._process_docx(request.content, request.custom_params)
            elif document_type == 'xlsx':
                extracted_content = await self._process_xlsx(request.content, request.custom_params)
            elif document_type == 'pptx':
                extracted_content = await self._process_pptx(request.content, request.custom_params)
            else:
                extracted_content = await self._process_generic_office(request.content, document_type)

            # Create chunks
            chunks = await self._create_chunks(extracted_content, request)

            # Generate metadata
            metadata = await self._extract_metadata(extracted_content, request, document_type)

            # Calculate quality metrics
            await self._calculate_quality_metrics(quality_session, extracted_content, chunks)

            response = ParseResponse(
                content=extracted_content['text'],
                chunks=chunks,
                metadata=metadata,
                format=document_format,
                processing_time=(datetime.now() - processing_start).total_seconds(),
                success=True
            )

            # End quality monitoring
            self.quality_monitor.end_session(session_id, success=True)

            return response

        except Exception as e:
            # End quality monitoring with error
            self.quality_monitor.end_session(session_id, success=False, error=e)

            return ParseResponse(
                content="",
                chunks=[],
                metadata=Metadata(),
                format=DocumentFormat.UNKNOWN,
                processing_time=(datetime.now() - processing_start).total_seconds(),
                success=False,
                error=str(e)
            )

    def _detect_office_format(self, request: ParseRequest) -> DocumentFormat:
        """Detect office document format."""
        if request.file_path:
            file_ext = Path(request.file_path).suffix.lower()
            if file_ext == '.pdf':
                return DocumentFormat.PDF
            elif file_ext in ['.doc', '.docx']:
                return DocumentFormat.DOCX
            elif file_ext in ['.xls', '.xlsx']:
                return DocumentFormat.XLSX
            elif file_ext in ['.ppt', '.pptx']:
                return DocumentFormat.PPTX

        # Check MIME type
        if request.mime_type:
            if 'pdf' in request.mime_type:
                return DocumentFormat.PDF
            elif 'word' in request.mime_type or 'docx' in request.mime_type:
                return DocumentFormat.DOCX
            elif 'sheet' in request.mime_type or 'excel' in request.mime_type:
                return DocumentFormat.XLSX
            elif 'presentation' in request.mime_type or 'powerpoint' in request.mime_type:
                return DocumentFormat.PPTX

        # Check content magic bytes
        if request.content:
            content_prefix = request.content[:16]
            if content_prefix.startswith(b'%PDF-'):
                return DocumentFormat.PDF
            elif content_prefix.startswith(b'PK\x03\x04'):  # ZIP-based (DOCX, XLSX, PPTX)
                return DocumentFormat.DOCX  # Default for ZIP-based

        return DocumentFormat.UNKNOWN

    async def _process_pdf(self, content: bytes, params: Optional[Dict[str, Any]]) -> Dict[str, Any]:
        """Process PDF content."""
        if not content:
            return {'text': '', 'features': OfficeDocumentFeatures('pdf', 0, False, False, False, 0, 0, 0, False, False, False)}

        text = ""
        page_count = 0
        has_images = False
        has_tables = False

        try:
            # Try PyMuPDF first (better performance)
            if PYMUPDF_AVAILABLE:
                pdf_document = fitz.open(stream=content, filetype="pdf")
                page_count = pdf_document.page_count

                for page_num in range(page_count):
                    page = pdf_document[page_num]

                    # Extract text
                    page_text = page.get_text()
                    text += page_text + "\n"

                    # Check for images
                    if not has_images and page.get_images():
                        has_images = True

                    # Check for tables (basic detection)
                    if not has_tables and '|' in page_text and '\n' in page_text:
                        lines = page_text.split('\n')
                        for line in lines[:10]:  # Check first 10 lines
                            if '|' in line and line.count('|') >= 3:
                                has_tables = True
                                break

                pdf_document.close()

            # Fallback to PyPDF2
            elif PYPDF2_AVAILABLE:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                page_count = len(pdf_reader.pages)

                for page in pdf_reader.pages:
                    try:
                        page_text = page.extract_text()
                        text += page_text + "\n"
                    except Exception:
                        continue

        except Exception as e:
            raise ValueError(f"Failed to process PDF: {e}")

        return {
            'text': text,
            'features': OfficeDocumentFeatures(
                document_type='pdf',
                page_count=page_count,
                has_images=has_images,
                has_tables=has_tables,
                has_formulas=False,
                slide_count=0,
                worksheet_count=0,
                section_count=page_count,
                is_protected=False,
                has_comments=False,
                has_track_changes=False
            )
        }

    async def _process_docx(self, content: bytes, params: Optional[Dict[str, Any]]) -> Dict[str, Any]:
        """Process DOCX content."""
        if not DOCX_AVAILABLE:
            return await self._process_fallback_text(content)

        try:
            doc = DocxDocument(io.BytesIO(content))
            text = ""
            has_tables = False
            has_images = False
            section_count = len(doc.sections)

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"

            # Extract tables
            tables = []
            for table in doc.tables:
                has_tables = True
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    if any(row_data):  # Skip empty rows
                        table_data.append(row_data)
                if table_data:
                    tables.append(table_data)
                    # Add table content to main text
                    text += "\n".join([" | ".join(row) for row in table_data]) + "\n"

            # Check for images (basic detection)
            try:
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        has_images = True
                        break
            except:
                pass

            return {
                'text': text,
                'tables': tables,
                'features': OfficeDocumentFeatures(
                    document_type='docx',
                    page_count=section_count,
                    has_images=has_images,
                    has_tables=has_tables,
                    has_formulas=False,
                    slide_count=0,
                    worksheet_count=0,
                    section_count=section_count,
                    is_protected=False,
                    has_comments=False,
                    has_track_changes=False
                )
            }

        except Exception as e:
            raise ValueError(f"Failed to process DOCX: {e}")

    async def _process_xlsx(self, content: bytes, params: Optional[Dict[str, Any]]) -> Dict[str, Any]:
        """Process XLSX content."""
        if not XLSX_AVAILABLE:
            return await self._process_fallback_text(content)

        try:
            workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            text = ""
            has_tables = True
            has_formulas = False
            worksheet_count = len(workbook.worksheets)

            tables = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []

                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        # Convert all cells to string
                        row_text = [str(cell) if cell is not None else "" for cell in row]
                        sheet_data.append(row_text)

                        # Add to main text
                        text += " | ".join(row_text) + "\n"

                if sheet_data:
                    tables.append({
                        'sheet_name': sheet_name,
                        'data': sheet_data
                    })

            workbook.close()

            return {
                'text': text,
                'tables': tables,
                'features': OfficeDocumentFeatures(
                    document_type='xlsx',
                    page_count=0,
                    has_images=False,
                    has_tables=has_tables,
                    has_formulas=has_formulas,
                    slide_count=0,
                    worksheet_count=worksheet_count,
                    section_count=worksheet_count,
                    is_protected=False,
                    has_comments=False,
                    has_track_changes=False
                )
            }

        except Exception as e:
            raise ValueError(f"Failed to process XLSX: {e}")

    async def _process_pptx(self, content: bytes, params: Optional[Dict[str, Any]]) -> Dict[str, Any]:
        """Process PPTX content."""
        if not PPTX_AVAILABLE:
            return await self._process_fallback_text(content)

        try:
            prs = Presentation(io.BytesIO(content))
            text = ""
            has_images = False
            slide_count = len(prs.slides)

            for slide in prs.slides:
                slide_text = ""

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"

                    # Check for images
                    if not has_images and hasattr(shape, "image"):
                        has_images = True

                text += slide_text + "\n"

            return {
                'text': text,
                'features': OfficeDocumentFeatures(
                    document_type='pptx',
                    page_count=slide_count,
                    has_images=has_images,
                    has_tables=False,
                    has_formulas=False,
                    slide_count=slide_count,
                    worksheet_count=0,
                    section_count=slide_count,
                    is_protected=False,
                    has_comments=False,
                    has_track_changes=False
                )
            }

        except Exception as e:
            raise ValueError(f"Failed to process PPTX: {e}")

    async def _process_generic_office(self, content: bytes, document_type: str) -> Dict[str, Any]:
        """Process generic office document."""
        return await self._process_fallback_text(content)

    async def _process_fallback_text(self, content: bytes) -> Dict[str, Any]:
        """Fallback text processing for unsupported formats."""
        try:
            text = content.decode('utf-8', errors='ignore')
            return {
                'text': text,
                'features': OfficeDocumentFeatures(
                    document_type='unknown',
                    page_count=1,
                    has_images=False,
                    has_tables=False,
                    has_formulas=False,
                    slide_count=0,
                    worksheet_count=0,
                    section_count=1,
                    is_protected=False,
                    has_comments=False,
                    has_track_changes=False
                )
            }
        except Exception:
            return {
                'text': '',
                'features': OfficeDocumentFeatures(
                    document_type='unknown',
                    page_count=0,
                    has_images=False,
                    has_tables=False,
                    has_formulas=False,
                    slide_count=0,
                    worksheet_count=0,
                    section_count=0,
                    is_protected=False,
                    has_comments=False,
                    has_track_changes=False
                )
            }

    async def _create_chunks(self, content: Dict[str, Any], request: ParseRequest) -> List[DocumentChunk]:
        """Create document chunks from processed content."""
        text = content['text']
        if not text.strip():
            return []

        # Get chunking parameters
        params = request.custom_params or {}
        chunk_size = params.get('chunk_size', 800)
        overlap = params.get('overlap_size', 200)

        # Create chunks
        chunks = []
        words = text.split()
        current_chunk = []
        current_length = 0

        for i, word in enumerate(words):
            if current_length + len(word) + 1 > chunk_size and current_chunk:
                # Save current chunk
                chunk_text = ' '.join(current_chunk)
                chunks.append(DocumentChunk(
                    content=chunk_text,
                    metadata=Metadata({
                        'chunk_index': len(chunks),
                        'word_count': len(current_chunk),
                        'source': 'office_documents',
                        'document_type': content['features'].document_type
                    })
                ))

                # Start new chunk with overlap
                overlap_words = min(overlap // 6, len(current_chunk))
                current_chunk = current_chunk[-overlap_words:] if overlap_words > 0 else []
                current_length = sum(len(w) + 1 for w in current_chunk)

            current_chunk.append(word)
            current_length += len(word) + 1

        # Add final chunk
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            chunks.append(DocumentChunk(
                content=chunk_text,
                metadata=Metadata({
                    'chunk_index': len(chunks),
                    'word_count': len(current_chunk),
                    'source': 'office_documents',
                    'document_type': content['features'].document_type
                })
            ))

        return chunks

    async def _extract_metadata(
        self,
        content: Dict[str, Any],
        request: ParseRequest,
        document_type: str
    ) -> Metadata:
        """Extract metadata from processed content."""
        features = content['features']

        metadata_dict = {
            'source_type': 'office_documents',
            'document_type': document_type,
            'processing_strategy': request.strategy.value,
            'page_count': features.page_count,
            'has_images': features.has_images,
            'has_tables': features.has_tables,
            'slide_count': features.slide_count,
            'worksheet_count': features.worksheet_count,
            'section_count': features.section_count,
            'is_protected': features.is_protected,
            'has_comments': features.has_comments,
            'has_track_changes': features.has_track_changes
        }

        # Add file path if available
        if request.file_path:
            metadata_dict['file_path'] = request.file_path

        # Add tables info
        if 'tables' in content and content['tables']:
            metadata_dict['table_count'] = len(content['tables'])

        return Metadata(metadata_dict)

    async def _calculate_quality_metrics(
        self,
        session_id: str,
        content: Dict[str, Any],
        chunks: List[DocumentChunk]
    ):
        """Calculate quality metrics for the processing session."""
        # Text extraction quality
        text = content['text']
        if text.strip():
            text_quality = min(1.0, len(text.strip()) / 500)  # More generous for office docs
            self.quality_monitor.add_measurement(
                session_id, QualityMetric.TEXT_EXTRACTION, text_quality
            )

        # Structure preservation
        features = content['features']
        structure_score = 0.0
        if features.has_tables:
            structure_score += 0.3
        if features.has_images:
            structure_score += 0.2
        if features.page_count > 1:
            structure_score += 0.2
        if features.section_count > 1:
            structure_score += 0.2
        if content.get('tables'):
            structure_score += 0.1
        structure_score = min(1.0, structure_score)

        self.quality_monitor.add_measurement(
            session_id, QualityMetric.STRUCTURE_PRESERVATION, structure_score
        )

        # Content completeness
        completeness_score = min(1.0, len(text) / 200)  # Minimum 200 chars for office docs
        self.quality_monitor.add_measurement(
            session_id, QualityMetric.CONTENT_COMPLETENESS, completeness_score
        )