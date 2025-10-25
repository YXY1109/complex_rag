"""
办公文档处理器

此模块实现PDF、DOCX、Excel、PPT等办公文档的专用处理器，
参考RAGFlow deepdoc模块中的文档处理逻辑。
"""

import asyncio
import os
import tempfile
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from docx import Document
import zipfile
import xml.etree.ElementTree as ET
from typing import Dict, Any, List, Optional, Union, BinaryIO
from pathlib import Path
import io
import uuid

from .base_processor import BaseProcessor
from ..interfaces.parser_interface import (
    ParseResult,
    DocumentMetadata,
    DocumentType,
    ProcessingStrategy,
    TextChunk,
    ImageInfo,
    TableInfo,
    ParseException,
    UnsupportedFormatError
)
from ..strategy_config import ProcessingStrategyConfig


class OfficeDocumentProcessor(BaseProcessor):
    """
    办公文档处理器。

    专门处理PDF、DOCX、Excel、PPT等办公格式文档。
    """

    def __init__(self, config):
        """
        初始化办公文档处理器。

        Args:
            config: 解析器配置
        """
        super().__init__(config)
        self.supported_extensions = {
            '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'
        }

    async def initialize(self) -> bool:
        """
        初始化处理器。

        Returns:
            bool: 初始化是否成功
        """
        try:
            # 测试依赖库
            import fitz  # PyMuPDF
            import pandas
            import pptx
            import docx
            return True
        except ImportError as e:
            print(f"办公文档处理器初始化失败，缺少依赖: {e}")
            return False

    async def cleanup(self) -> None:
        """清理处理器资源。"""
        pass

    async def _parse_with_config(
        self,
        file_path: str,
        strategy: ProcessingStrategy,
        config: ProcessingStrategyConfig,
        **kwargs
    ) -> ParseResult:
        """
        使用配置解析办公文档。

        Args:
            file_path: 文件路径
            strategy: 处理策略
            config: 策略配置
            **kwargs: 额外参数

        Returns:
            ParseResult: 解析结果
        """
        # 验证文件
        self._validate_file(file_path)

        # 确定文档类型
        file_ext = Path(file_path).suffix.lower()

        try:
            if file_ext == '.pdf':
                return await self._parse_pdf_document(file_path, strategy, config)
            elif file_ext in {'.doc', '.docx'}:
                return await self._parse_word_document(file_path, strategy, config)
            elif file_ext in {'.xls', '.xlsx'}:
                return await self._parse_excel_document(file_path, strategy, config)
            elif file_ext in {'.ppt', '.pptx'}:
                return await self._parse_powerpoint_document(file_path, strategy, config)
            else:
                raise UnsupportedFormatError(
                    f"不支持的办公文档格式: {file_ext}",
                    parser=self.parser_name,
                    file_path=file_path
                )
        except Exception as e:
            raise ParseException(
                f"办公文档解析失败: {str(e)}",
                parser=self.parser_name,
                file_path=file_path
            )

    async def _parse_pdf_document(
        self,
        file_path: str,
        strategy: ProcessingStrategy,
        config: ProcessingStrategyConfig
    ) -> ParseResult:
        """
        解析PDF文档。

        Args:
            file_path: 文件路径
            strategy: 处理策略
            config: 策略配置

        Returns:
            ParseResult: 解析结果
        """
        try:
            doc = fitz.open(file_path)
            page_count = doc.page_count

            # 提取元数据
            metadata = await self._extract_pdf_metadata(doc, file_path, page_count)

            # 根据策略选择处理方式
            if strategy == ProcessingStrategy.MULTIMODAL_ANALYSIS:
                result = await self._parse_pdf_multimodal(doc, config, metadata)
            elif strategy == ProcessingStrategy.TABLE_EXTRACTION:
                result = await self._parse_pdf_with_tables(doc, config, metadata)
            elif strategy == ProcessingStrategy.PRESERVE_LAYOUT:
                result = await self._parse_pdf_with_layout(doc, config, metadata)
            else:
                result = await self._parse_pdf_basic(doc, config, metadata)

            doc.close()
            return result

        except Exception as e:
            raise ParseException(f"PDF解析失败: {str(e)}", parser=self.parser_name, file_path=file_path)

    async def _parse_pdf_basic(
        self,
        doc,
        config: ProcessingStrategyConfig,
        metadata: DocumentMetadata
    ) -> ParseResult:
        """
        基础PDF解析。

        Args:
            doc: PDF文档对象
            config: 配置
            metadata: 元数据

        Returns:
            ParseResult: 解析结果
        """
        full_text = ""
        text_chunks = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            page_text = page.get_text()

            if page_text.strip():
                full_text += page_text + "\n"

                # 创建页级文本块
                chunk = TextChunk(
                    content=page_text.strip(),
                    page_number=page_num + 1,
                    chunk_id=f"page_{page_num + 1}",
                    confidence=0.9
                )
                text_chunks.append(chunk)

        # 更新元数据
        metadata.word_count = len(full_text.split())
        metadata.character_count = len(full_text)

        return ParseResult(
            success=True,
            metadata=metadata,
            full_text=full_text.strip(),
            text_chunks=text_chunks
        )

    async def _parse_pdf_with_layout(
        self,
        doc,
        config: ProcessingStrategyConfig,
        metadata: DocumentMetadata
    ) -> ParseResult:
        """
        布局保持的PDF解析。

        Args:
            doc: PDF文档对象
            config: 配置
            metadata: 元数据

        Returns:
            ParseResult: 解析结果
        """
        full_text = ""
        text_chunks = []
        images = []

        for page_num in range(len(doc)):
            page = doc[page_num]

            # 提取文本块
            blocks = page.get_text("dict")
            page_text = page.get_text()
            full_text += page_text + "\n"

            # 处理文本块
            if "blocks" in blocks:
                for block in blocks["blocks"]:
                    if block["type"] == 0:  # 文本块
                        chunk = TextChunk(
                            content=block["text"],
                            page_number=page_num + 1,
                            chunk_id=f"block_{page_num + 1}_{len(text_chunks)}",
                            bbox=block.get("bbox"),
                            confidence=0.85
                        )
                        text_chunks.append(chunk)

            # 提取图像
            if config.image.extract_images:
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n < 5:  # 跳过小图像
                            pix = None
                            continue

                        image_info = ImageInfo(
                            image_id=f"img_{page_num + 1}_{img_index}",
                            page_number=page_num + 1,
                            bbox=[0, 0, pix.width, pix.height],
                            width=pix.width,
                            height=pix.height,
                            format="png",
                            metadata={
                                "xref": xref,
                                "colorspace": pix.colorspace.name if pix.colorspace else "unknown"
                            }
                        )
                        images.append(image_info)
                        pix = None
                    except Exception:
                        continue

        # 创建解析结果
        result = ParseResult(
            success=True,
            metadata=metadata,
            full_text=full_text.strip(),
            text_chunks=text_chunks,
            images=images,
            structured_data={
                "layout_blocks": len(text_chunks),
                "image_count": len(images)
            }
        )

        return result

    async def _parse_pdf_multimodal(
        self,
        doc,
        config: ProcessingStrategyConfig,
        metadata: DocumentMetadata
    ) -> ParseResult:
        """
        多模态PDF解析。

        Args:
            doc: PDF文档对象
            config: 配置
            metadata: 元数据

        Returns:
            ParseResult: 解析结果
        """
        # 这里实现更复杂的多模态分析
        # 包括OCR、图像分析、表格提取等

        # 先进行基础解析
        basic_result = await self._parse_pdf_with_layout(doc, config, metadata)

        # 如果启用了OCR，处理扫描页面
        if config.ocr.enabled:
            # 这里可以集成OCR引擎
            # 由于篇幅限制，这里只是框架
            pass

        # 如果启用了表格提取
        if config.table.enabled:
            tables = await self._extract_pdf_tables(doc, config)
            basic_result.tables = tables

        return basic_result

    async def _parse_pdf_with_tables(
        self,
        doc,
        config: ProcessingStrategyConfig,
        metadata: DocumentMetadata
    ) -> ParseResult:
        """
        带表格提取的PDF解析。

        Args:
            doc: PDF文档对象
            config: 配置
            metadata: 元数据

        Returns:
            ParseResult: 解析结果
        """
        # 先进行基础解析
        basic_result = await self._parse_pdf_basic(doc, config, metadata)

        # 提取表格
        tables = await self._extract_pdf_tables(doc, config)
        basic_result.tables = tables
        basic_result.metadata.has_tables = len(tables) > 0

        return basic_result

    async def _extract_pdf_tables(
        self,
        doc,
        config: ProcessingStrategyConfig
    ) -> List[TableInfo]:
        """
        从PDF提取表格。

        Args:
            doc: PDF文档对象
            config: 配置

        Returns:
            List[TableInfo]: 表格信息列表
        """
        tables = []

        try:
            for page_num in range(len(doc)):
                page = doc[page_num]
                # 使用camelot或tabula进行表格提取
                # 这里简化处理，实际应用中需要集成专门的表格提取库
                tables_found = page.find_tables()

                for i, table in enumerate(tables_found):
                    table_data = table.extract()
                    if table_data:
                        # 转换表格数据
                        rows = len(table_data)
                        cols = len(table_data[0]) if table_data else 0

                        # 提取表头
                        headers = table_data[0] if rows > 0 else []

                        table_info = TableInfo(
                            table_id=f"table_{page_num + 1}_{i}",
                            page_number=page_num + 1,
                            bbox=table.bbox,
                            rows=rows,
                            columns=cols,
                            headers=headers,
                            data=table_data,
                            confidence=0.8
                        )
                        tables.append(table_info)
        except Exception as e:
            print(f"PDF表格提取警告: {e}")

        return tables

    async def _extract_pdf_metadata(
        self,
        doc,
        file_path: str,
        page_count: int
    ) -> DocumentMetadata:
        """
        提取PDF元数据。

        Args:
            doc: PDF文档对象
            file_path: 文件路径
            page_count: 页数

        Returns:
            DocumentMetadata: 元数据
        """
        file_size = os.path.getsize(file_path)
        file_name = os.path.basename(file_path)

        # 提取PDF元数据
        metadata_dict = doc.metadata
        title = metadata_dict.get('title', file_name)
        author = metadata_dict.get('author')
        creator = metadata_dict.get('creator')
        producer = metadata_dict.get('producer')
        creation_date = metadata_dict.get('creationDate')
        mod_date = metadata_dict.get('modDate')

        # 检查加密状态
        encrypted = doc.is_encrypted

        # 检查图像和表格
        has_images = False
        has_tables = False
        for page in doc:
            if page.get_images():
                has_images = True
            if page.find_tables():
                has_tables = True
            if has_images and has_tables:
                break

        return DocumentMetadata(
            file_name=file_name,
            file_size=file_size,
            file_type=DocumentType.PDF,
            mime_type="application/pdf",
            title=title,
            author=author,
            created_date=str(creation_date) if creation_date else None,
            modified_date=str(mod_date) if mod_date else None,
            page_count=page_count,
            encryption=encrypted,
            has_images=has_images,
            has_tables=has_tables,
            metadata={
                "creator": creator,
                "producer": producer,
                "pdf_version": metadata_dict.get('pdfVersion'),
                "keywords": metadata_dict.get('keywords')
            }
        )

    async def _parse_word_document(
        self,
        file_path: str,
        strategy: ProcessingStrategy,
        config: ProcessingStrategyConfig
    ) -> ParseResult:
        """
        解析Word文档。

        Args:
            file_path: 文件路径
            strategy: 处理策略
            config: 配置

        Returns:
            ParseResult: 解析结果
        """
        try:
            doc = Document(file_path)
            full_text = ""
            text_chunks = []
            images = []

            # 提取段落文本
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    full_text += paragraph.text + "\n"

                    chunk = TextChunk(
                        content=paragraph.text.strip(),
                        chunk_id=f"para_{i}",
                        confidence=0.9
                    )
                    text_chunks.append(chunk)

            # 提取表格
            tables = []
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)

                if table_data:
                    table_info = TableInfo(
                        table_id=f"table_{i}",
                        page_number=1,
                        bbox=[0, 0, 0, 0],
                        rows=len(table_data),
                        columns=len(table_data[0]) if table_data else 0,
                        headers=table_data[0] if table_data else [],
                        data=table_data,
                        confidence=0.9
                    )
                    tables.append(table_info)

            # 提取图像（DOCX格式）
            if file_path.endswith('.docx'):
                # 这里可以提取DOCX中的图像
                # 由于篇幅限制，这里只是框架
                pass

            # 创建元数据
            file_size = os.path.getsize(file_path)
            file_name = os.path.basename(file_path)

            metadata = DocumentMetadata(
                file_name=file_name,
                file_size=file_size,
                file_type=DocumentType.DOCX if file_path.endswith('.docx') else DocumentType.DOC,
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document" if file_path.endswith('.docx') else "application/msword",
                title=doc.core_properties.title or file_name,
                author=doc.core_properties.author,
                created_date=str(doc.core_properties.created) if doc.core_properties.created else None,
                modified_date=str(doc.core_properties.modified) if doc.core_properties.modified else None,
                page_count=1,
                word_count=len(full_text.split()),
                character_count=len(full_text),
                has_tables=len(tables) > 0,
                metadata={
                    "subject": doc.core_properties.subject,
                    "keywords": doc.core_properties.keywords,
                    "paragraphs": len(doc.paragraphs),
                    "tables": len(tables)
                }
            )

            return ParseResult(
                success=True,
                metadata=metadata,
                full_text=full_text.strip(),
                text_chunks=text_chunks,
                tables=tables,
                images=images,
                structured_data={
                    "document_stats": {
                        "paragraphs": len(doc.paragraphs),
                        "sections": len(doc.sections),
                        "tables": len(tables)
                    }
                }
            )

        except Exception as e:
            raise ParseException(f"Word文档解析失败: {str(e)}", parser=self.parser_name, file_path=file_path)

    async def _parse_excel_document(
        self,
        file_path: str,
        strategy: ProcessingStrategy,
        config: ProcessingStrategyConfig
    ) -> ParseResult:
        """
        解析Excel文档。

        Args:
            file_path: 文件路径
            strategy: 处理策略
            config: 配置

        Returns:
            ParseResult: 解析结果
        """
        try:
            # 读取Excel文件
            excel_file = pd.ExcelFile(file_path)
            sheets = excel_file.sheet_names

            full_text = ""
            tables = []

            # 处理每个工作表
            for sheet_name in sheets:
                df = pd.read_excel(file_path, sheet_name=sheet_name)

                # 转换为文本
                sheet_text = f"工作表: {sheet_name}\n"
                sheet_text += df.to_string(index=False)
                full_text += sheet_text + "\n\n"

                # 创建表格信息
                table_data = df.values.tolist()
                headers = df.columns.tolist()

                table_info = TableInfo(
                    table_id=f"sheet_{sheet_name}",
                    page_number=1,
                    bbox=[0, 0, 0, 0],
                    rows=len(table_data),
                    columns=len(headers),
                    headers=headers,
                    data=table_data,
                    confidence=0.95
                )
                tables.append(table_info)

            # 创建元数据
            file_size = os.path.getsize(file_path)
            file_name = os.path.basename(file_path)

            metadata = DocumentMetadata(
                file_name=file_name,
                file_size=file_size,
                file_type=DocumentType.XLSX if file_path.endswith('.xlsx') else DocumentType.XLS,
                mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" if file_path.endswith('.xlsx') else "application/vnd.ms-excel",
                title=file_name,
                page_count=1,
                has_tables=True,
                metadata={
                    "sheets": sheets,
                    "sheet_count": len(sheets)
                }
            )

            # 创建文本块（按工作表分块）
            text_chunks = []
            for i, sheet_name in enumerate(sheets):
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheet_text = f"工作表: {sheet_name}\n" + df.to_string(index=False)

                chunk = TextChunk(
                    content=sheet_text,
                    chunk_id=f"sheet_{i}",
                    confidence=0.95
                )
                text_chunks.append(chunk)

            return ParseResult(
                success=True,
                metadata=metadata,
                full_text=full_text.strip(),
                text_chunks=text_chunks,
                tables=tables,
                structured_data={
                    "sheets": sheets,
                    "sheet_count": len(sheets)
                }
            )

        except Exception as e:
            raise ParseException(f"Excel文档解析失败: {str(e)}", parser=self.parser_name, file_path=file_path)

    async def _parse_powerpoint_document(
        self,
        file_path: str,
        strategy: ProcessingStrategy,
        config: ProcessingStrategyConfig
    ) -> ParseResult:
        """
        解析PowerPoint文档。

        Args:
            file_path: 文件路径
            strategy: 处理策略
            config: 配置

        Returns:
            ParseResult: 解析结果
        """
        try:
            prs = Presentation(file_path)
            full_text = ""
            text_chunks = []
            images = []

            # 提取每页内容
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"幻灯片 {slide_num + 1}:\n"

                # 提取文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"

                full_text += slide_text + "\n"

                # 创建幻灯片文本块
                chunk = TextChunk(
                    content=slide_text.strip(),
                    page_number=slide_num + 1,
                    chunk_id=f"slide_{slide_num + 1}",
                    confidence=0.9
                )
                text_chunks.append(chunk)

                # 提取图像
                if config.image.extract_images:
                    for shape in slide.shapes:
                        if hasattr(shape, "image"):
                            # 这里可以提取图像
                            # 由于篇幅限制，这里只是框架
                            pass

            # 创建元数据
            file_size = os.path.getsize(file_path)
            file_name = os.path.basename(file_path)

            metadata = DocumentMetadata(
                file_name=file_name,
                file_size=file_size,
                file_type=DocumentType.PPTX if file_path.endswith('.pptx') else DocumentType.PPT,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation" if file_path.endswith('.pptx') else "application/vnd.ms-powerpoint",
                title=file_name,
                page_count=len(prs.slides),
                has_images=False,  # 这里简化处理
                metadata={
                    "slide_count": len(prs.slides)
                }
            )

            return ParseResult(
                success=True,
                metadata=metadata,
                full_text=full_text.strip(),
                text_chunks=text_chunks,
                images=images,
                structured_data={
                    "slide_count": len(prs.slides)
                }
            )

        except Exception as e:
            raise ParseException(f"PowerPoint文档解析失败: {str(e)}", parser=self.parser_name, file_path=file_path)


# 导出
__all__ = ['OfficeDocumentProcessor']